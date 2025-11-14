import os
import json
import uuid
import datetime
from functools import wraps
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import jwt
from werkzeug.utils import secure_filename
from io import StringIO
import PyPDF2
import requests
import traceback

try:
    # Try loading .env from the backend directory if python-dotenv is installed
    from dotenv import load_dotenv
    load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))
except Exception:
    # dotenv not installed or .env missing — that's fine, env vars will be read from the environment
    pass

try:
    import bcrypt
except ImportError:
    bcrypt = None

try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    PPTXPresentation = None

SECRET_KEY = os.environ.get('SECRET_KEY', 'dev-secret')
DATA_FILE = os.path.join(os.path.dirname(__file__), 'data.json')
UPLOAD_DIR = os.path.join(os.path.dirname(__file__), 'uploads')

ALLOWED_TEXT_EXTS = {'txt', 'pdf', 'doc', 'docx'}
ALLOWED_IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'gif'}

os.makedirs(UPLOAD_DIR, exist_ok=True)

app = Flask(__name__)
CORS(app)
app.config['UPLOAD_FOLDER'] = UPLOAD_DIR
app.config['SECRET_KEY'] = SECRET_KEY
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload

# Simple helpers to read/write a JSON datastore (file-backed)

def read_data():
    try:
        with open(DATA_FILE, 'r') as f:
            return json.load(f)
    except Exception:
        return {'users': {}, 'presentations': {}}


def write_data(d):
    with open(DATA_FILE, 'w') as f:
        json.dump(d, f, indent=2)


def hash_password(password):
    """Hash password with bcrypt or fallback to plaintext with warning"""
    if bcrypt:
        return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
    else:
        print("WARNING: bcrypt not installed, storing plaintext passwords!")
        return password


def verify_password(password, hashed):
    """Verify password with bcrypt or fallback to plaintext comparison"""
    if bcrypt:
        return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))
    else:
        return password == hashed


def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = None
        if 'Authorization' in request.headers:
            parts = request.headers.get('Authorization').split()
            if len(parts) == 2 and parts[0].lower() == 'bearer':
                token = parts[1]
        if not token:
            return jsonify({'message': 'Token is missing!'}), 401
        try:
            data = jwt.decode(token, app.config['SECRET_KEY'], algorithms=['HS256'])
            request.user = data['username']
        except Exception as e:
            return jsonify({'message': 'Token is invalid!', 'error': str(e)}), 401
        return f(*args, **kwargs)
    return decorated


# ============== AUTHENTICATION ==============

@app.route('/signup', methods=['POST'])
def signup():
    try:
        payload = request.json or {}
        username = payload.get('username', '').strip()
        password = payload.get('password', '').strip()
        
        if not username or len(username) < 3:
            return jsonify({'message': 'Username must be at least 3 characters'}), 400
        if not password or len(password) < 6:
            return jsonify({'message': 'Password must be at least 6 characters'}), 400
        
        data = read_data()
        users = data['users']
        
        if username in users:
            return jsonify({'message': 'User already exists'}), 400
        
        users[username] = {
            'password': hash_password(password),
            'created_at': datetime.datetime.utcnow().isoformat()
        }
        write_data(data)
        return jsonify({'message': 'User created successfully'}), 201
    except Exception as e:
        return jsonify({'message': 'Signup failed', 'error': str(e)}), 500


@app.route('/login', methods=['POST'])
def login():
    try:
        payload = request.json or {}
        username = payload.get('username', '').strip()
        password = payload.get('password', '')
        
        if not username or not password:
            return jsonify({'message': 'Username and password required'}), 400
        
        data = read_data()
        users = data['users']
        
        if username not in users or not verify_password(password, users[username]['password']):
            return jsonify({'message': 'Invalid credentials'}), 401
        
        token = jwt.encode(
            {'username': username, 'exp': datetime.datetime.utcnow() + datetime.timedelta(hours=12)},
            app.config['SECRET_KEY'],
            algorithm='HS256'
        )
        return jsonify({'token': token, 'username': username}), 200
    except Exception as e:
        return jsonify({'message': 'Login failed', 'error': str(e)}), 500


# ============== PRESENTATIONS ==============

@app.route('/presentations', methods=['GET'])
@token_required
def list_presentations():
    try:
        username = request.user
        data = read_data()
        pres = data.get('presentations', {})
        user_pres = [v for v in pres.values() if v['owner'] == username]
        return jsonify({'presentations': user_pres}), 200
    except Exception as e:
        return jsonify({'message': 'Failed to list presentations', 'error': str(e)}), 500


@app.route('/presentations', methods=['POST'])
@token_required
def create_presentation():
    try:
        payload = request.json or {}
        username = request.user
        title = payload.get('title', 'Untitled').strip()
        
        if not title:
            return jsonify({'message': 'Title is required'}), 400
        
        pres_id = str(uuid.uuid4())
        data = read_data()
        slide_count = max(1, min(int(payload.get('slide_count', 5)), 50))  # 1-50 slides
        
        slides = []
        for i in range(slide_count):
            slides.append({
                'id': str(uuid.uuid4()),
                'title': f'Slide {i+1}',
                'content': 'Add your content here',
                'image': None,
                'style': {
                    'titleFontSize': 32,
                    'contentFontSize': 18,
                    'fontColor': '#000000',
                    'backgroundColor': '#ffffff',
                    'backgroundImage': None,
                    'backgroundOpacity': 100,
                    'backgroundBlur': 0
                }
            })
        
        data['presentations'][pres_id] = {
            'id': pres_id,
            'owner': username,
            'title': title,
            'slides': slides,
            'created_at': datetime.datetime.utcnow().isoformat(),
            'updated_at': datetime.datetime.utcnow().isoformat()
        }
        write_data(data)
        return jsonify({'presentation': data['presentations'][pres_id]}), 201
    except Exception as e:
        return jsonify({'message': 'Failed to create presentation', 'error': str(e)}), 500


@app.route('/presentations/<pres_id>', methods=['GET'])
@token_required
def get_presentation(pres_id):
    try:
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        return jsonify({'presentation': pres}), 200
    except Exception as e:
        return jsonify({'message': 'Failed to get presentation', 'error': str(e)}), 500


@app.route('/presentations/<pres_id>', methods=['PUT'])
@token_required
def update_presentation(pres_id):
    try:
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        payload = request.json or {}
        
        if 'title' in payload:
            title = payload['title'].strip()
            if not title:
                return jsonify({'message': 'Title cannot be empty'}), 400
            pres['title'] = title
        
        pres['updated_at'] = datetime.datetime.utcnow().isoformat()
        write_data(data)
        return jsonify({'presentation': pres}), 200
    except Exception as e:
        return jsonify({'message': 'Failed to update presentation', 'error': str(e)}), 500


@app.route('/presentations/<pres_id>', methods=['DELETE'])
@token_required
def delete_presentation(pres_id):
    try:
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        del data['presentations'][pres_id]
        write_data(data)
        return jsonify({'message': 'Presentation deleted'}), 200
    except Exception as e:
        return jsonify({'message': 'Failed to delete presentation', 'error': str(e)}), 500


# ============== SLIDES ==============

@app.route('/presentations/<pres_id>/slides', methods=['POST'])
@token_required
def create_slide(pres_id):
    try:
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        payload = request.json or {}
        slide = {
            'id': str(uuid.uuid4()),
            'title': payload.get('title', 'New Slide').strip(),
            'content': payload.get('content', ''),
            'image': payload.get('image'),
            'style': payload.get('style', {
                'titleFontSize': 32,
                'contentFontSize': 18,
                'fontColor': '#000000',
                'backgroundColor': '#ffffff',
                'backgroundImage': None,
                'backgroundOpacity': 100,
                'backgroundBlur': 0
            })
        }
        
        pres['slides'].append(slide)
        pres['updated_at'] = datetime.datetime.utcnow().isoformat()
        write_data(data)
        return jsonify({'slide': slide}), 201
    except Exception as e:
        return jsonify({'message': 'Failed to create slide', 'error': str(e)}), 500


@app.route('/presentations/<pres_id>/slides/<slide_id>', methods=['PUT'])
@token_required
def update_slide(pres_id, slide_id):
    try:
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        slide = next((s for s in pres['slides'] if s['id'] == slide_id), None)
        if not slide:
            return jsonify({'message': 'Slide not found'}), 404
        
        payload = request.json or {}
        
        # Update allowed fields
        if 'title' in payload:
            slide['title'] = payload['title'].strip()
        if 'content' in payload:
            slide['content'] = payload['content']
        if 'image' in payload:
            slide['image'] = payload['image']
        if 'style' in payload:
            slide['style'].update(payload['style'])
        
        pres['updated_at'] = datetime.datetime.utcnow().isoformat()
        write_data(data)
        return jsonify({'slide': slide}), 200
    except Exception as e:
        return jsonify({'message': 'Failed to update slide', 'error': str(e)}), 500


@app.route('/presentations/<pres_id>/slides/<slide_id>', methods=['DELETE'])
@token_required
def delete_slide(pres_id, slide_id):
    try:
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        pres['slides'] = [s for s in pres['slides'] if s['id'] != slide_id]
        pres['updated_at'] = datetime.datetime.utcnow().isoformat()
        write_data(data)
        return jsonify({'message': 'Slide deleted'}), 200
    except Exception as e:
        return jsonify({'message': 'Failed to delete slide', 'error': str(e)}), 500


@app.route('/presentations/<pres_id>/slides/reorder', methods=['POST'])
@token_required
def reorder_slides(pres_id):
    try:
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        payload = request.json or {}
        slide_order = payload.get('slide_ids', [])
        
        # Reorder slides based on provided IDs
        slide_map = {s['id']: s for s in pres['slides']}
        new_slides = []
        
        for sid in slide_order:
            if sid in slide_map:
                new_slides.append(slide_map[sid])
        
        # Add any slides that weren't in the order list (shouldn't happen)
        for s in pres['slides']:
            if s not in new_slides:
                new_slides.append(s)
        
        pres['slides'] = new_slides
        pres['updated_at'] = datetime.datetime.utcnow().isoformat()
        write_data(data)
        return jsonify({'presentation': pres}), 200
    except Exception as e:
        return jsonify({'message': 'Failed to reorder slides', 'error': str(e)}), 500


# ============== LLM HELPERS ==============

def call_llm_for_structured_content(prompt, context=''):
    """
    Call Anthropic (primary) or Gemini (fallback) with a structured prompt for deep analysis.
    Returns: dict with 'title' and 'bullets' (list of strings) or None if all fail.
    Falls back gracefully if no API keys are set.
    """
    full_prompt = f"""
    {prompt}

    Analyze the above deeply and structure your response as JSON:
    {{
    "title": "A concise, insightful title (max 10 words)",
    "bullets": [
        "First key insight or point (with reasoning if applicable)",
        "Second key insight or point",
        "Third key insight or point"
    ]
    }}

    Return ONLY valid JSON, no markdown or extra text.
    """

    # Try Anthropic first (using v1/messages endpoint for Claude 3.x)
    ANTHROPIC_KEY = os.environ.get('ANTHROPIC_API_KEY')
    if ANTHROPIC_KEY:
        try:
            anthropic_model = os.environ.get('ANTHROPIC_MODEL', 'claude-3-sonnet-20240229')
            anthropic_url = 'https://api.anthropic.com/v1/messages'
            anthropic_payload = {
                'model': anthropic_model,
                'max_tokens': int(os.environ.get('ANTHROPIC_MAX_TOKENS', 1024)),
                'messages': [
                    {'role': 'user', 'content': full_prompt}
                ]
            }
            resp = requests.post(
                anthropic_url,
                headers={
                    'x-api-key': ANTHROPIC_KEY,
                    'anthropic-version': '2023-06-01',
                    'Content-Type': 'application/json'
                },
                json=anthropic_payload,
                timeout=30,
            )
            if resp.status_code == 200:
                j = resp.json()
                # Extract text from content blocks
                text = ''
                if 'content' in j and isinstance(j['content'], list) and j['content']:
                    text = j['content'][0].get('text', '')
                if text:
                    # Try to extract JSON from the response
                    try:
                        # Clean the text to ensure it's a valid JSON object string
                        json_str = text.strip()
                        if json_str.startswith('```json'):
                            json_str = json_str[7:]
                        if json_str.endswith('```'):
                            json_str = json_str[:-3]
                        
                        result = json.loads(json_str.strip())
                        if 'title' in result and 'bullets' in result:
                            print(f'[Anthropic] Generated slide: {result["title"][:50]}...')
                            return result
                    except json.JSONDecodeError as je:
                        print(f'Anthropic response not JSON: {text[:100]}... Error: {je}')
            else:
                print(f'Anthropic returned {resp.status_code}: {resp.text[:200]}')
        except Exception as e:
            print(f'Anthropic call failed: {e}')
            traceback.print_exc()

    # Try Gemini/PaLM fallback
    GOOGLE_KEY = os.environ.get('GOOGLE_API_KEY') or os.environ.get('GOOGLE_KEY')
    if GOOGLE_KEY:
        try:
            GOOGLE_MODEL = os.environ.get('GOOGLE_MODEL', 'models/text-bison-001')
            gemini_url = f'https://generativelanguage.googleapis.com/v1beta2/{GOOGLE_MODEL}:generateText'
            params = {'key': GOOGLE_KEY}
            gemini_payload = {
                'prompt': {'text': full_prompt},
                'temperature': float(os.environ.get('GOOGLE_TEMPERATURE', 0.2)),
                'maxOutputTokens': int(os.environ.get('GOOGLE_MAX_TOKENS', 1024)),
            }
            resp = requests.post(gemini_url, params=params, json=gemini_payload, timeout=30)
            if resp.status_code in (200, 201):
                j = resp.json()
                text = ''
                if 'candidates' in j and j['candidates']:
                    text = j['candidates'][0].get('content', '')
                else:
                    text = j.get('output') or j.get('content') or ''
                if text:
                    try:
                        # Clean the text to ensure it's a valid JSON object string
                        json_str = text.strip()
                        if json_str.startswith('```json'):
                            json_str = json_str[7:]
                        if json_str.endswith('```'):
                            json_str = json_str[:-3]

                        result = json.loads(json_str.strip())
                        if 'title' in result and 'bullets' in result:
                            print(f'[Gemini] Generated slide: {result["title"][:50]}...')
                            return result
                    except json.JSONDecodeError as je:
                        print(f'Gemini response not JSON: {text[:100]}... Error: {je}')
            else:
                print(f'Gemini returned {resp.status_code}: {resp.text[:200]}')
        except Exception as e:
            print(f'Gemini call failed: {e}')

    # Fallback: return None (caller handles fallback)
    return None


# ============== FILE UPLOAD & PARSING ==============

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    try:
        text = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return None


def extract_text_from_document(file_path):
    """Extract text from TXT or DOC/DOCX files"""
    try:
        # Try reading as plaintext first
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"Text extraction error: {e}")
        return None


@app.route('/upload-document', methods=['POST'])
@token_required
def upload_document():
    try:
        if 'file' not in request.files:
            return jsonify({'message': 'No file provided'}), 400
        
        file = request.files['file']
        if not file or file.filename == '':
            return jsonify({'message': 'No selected file'}), 400
        
        # Validate file extension
        filename = secure_filename(file.filename)
        ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
        
        if ext not in ALLOWED_TEXT_EXTS:
            return jsonify({'message': f'Allowed formats: {", ".join(ALLOWED_TEXT_EXTS)}'}), 400
        
        # Save file temporarily
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{uuid.uuid4()}_{filename}")
        file.save(temp_path)
        
        # Extract text
        text = None
        if ext == 'pdf':
            text = extract_text_from_pdf(temp_path)
        else:
            text = extract_text_from_document(temp_path)
        
        # Clean up temp file
        try:
            os.remove(temp_path)
        except:
            pass
        
        if not text:
            return jsonify({'message': 'Could not extract text from document'}), 400
        
        return jsonify({'text': text, 'filename': filename}), 200
    except Exception as e:
        return jsonify({'message': 'Upload failed', 'error': str(e)}), 500


@app.route('/upload-image', methods=['POST'])
@token_required
def upload_image():
    try:
        if 'file' not in request.files:
            return jsonify({'message': 'No file provided'}), 400
        
        file = request.files['file']
        if not file or file.filename == '':
            return jsonify({'message': 'No selected file'}), 400
        
        # Validate file extension
        filename = secure_filename(file.filename)
        ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
        
        if ext not in ALLOWED_IMAGE_EXTS:
            return jsonify({'message': f'Allowed formats: {", ".join(ALLOWED_IMAGE_EXTS)}'}), 400
        
        # Save with unique name
        unique_filename = f"{uuid.uuid4()}_{filename}"
        save_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        file.save(save_path)
        
        url = f'/uploads/{unique_filename}'
        return jsonify({'url': url}), 201
    except Exception as e:
        return jsonify({'message': 'Upload failed', 'error': str(e)}), 500


@app.route('/uploads/<path:filename>')
def uploaded_file(filename):
    try:
        return send_from_directory(app.config['UPLOAD_FOLDER'], filename)
    except Exception as e:
        return jsonify({'message': 'File not found'}), 404


# ============== EXPORT ==============

@app.route('/presentations/<pres_id>/export', methods=['GET'])
@token_required
def export_presentation(pres_id):
    try:
        if not PPTXPresentation:
            return jsonify({'message': 'python-pptx not installed'}), 500
        
        data = read_data()
        pres = data.get('presentations', {}).get(pres_id)
        
        if not pres:
            return jsonify({'message': 'Presentation not found'}), 404
        if pres['owner'] != request.user:
            return jsonify({'message': 'Forbidden'}), 403
        
        # Create PPTX
        prs = PPTXPresentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for slide_data in pres['slides']:
            # Use blank slide layout
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            style = slide_data.get('style', {})
            bg_color = style.get('backgroundColor', '#ffffff')
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            
            # Parse hex color
            if bg_color.startswith('#'):
                r = int(bg_color[1:3], 16)
                g = int(bg_color[3:5], 16)
                b = int(bg_color[5:7], 16)
            else:
                r, g, b = 255, 255, 255
            
            fill.fore_color.rgb = RGBColor(r, g, b)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            p = title_frame.paragraphs[0]
            p.text = slide_data.get('title', '')
            p.font.size = Pt(style.get('titleFontSize', 32))
            p.font.bold = True
            
            font_color = style.get('fontColor', '#000000')
            if font_color.startswith('#'):
                r = int(font_color[1:3], 16)
                g = int(font_color[3:5], 16)
                b = int(font_color[5:7], 16)
            else:
                r, g, b = 0, 0, 0
            p.font.color.rgb = RGBColor(r, g, b)
            
            # Add content
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(4))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            p = content_frame.paragraphs[0]
            p.text = slide_data.get('content', '')
            p.font.size = Pt(style.get('contentFontSize', 18))
            p.font.color.rgb = RGBColor(r, g, b)
            
            # Add image if present
            if slide_data.get('image'):
                try:
                    img_path = slide_data['image']
                    if img_path.startswith('/uploads/'):
                        img_path = os.path.join(app.config['UPLOAD_FOLDER'], img_path.split('/')[-1])
                    
                    if os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, Inches(5.5), Inches(2.2), width=Inches(4))
                except Exception as e:
                    print(f"Could not add image: {e}")
        
        # Save to temp file
        output_file = os.path.join(app.config['UPLOAD_FOLDER'], f"{pres_id}_export.pptx")
        prs.save(output_file)
        
        # Return file
        return send_from_directory(
            app.config['UPLOAD_FOLDER'],
            os.path.basename(output_file),
            as_attachment=True,
            download_name=f"{pres['title']}.pptx"
        )
    except Exception as e:
        return jsonify({'message': 'Export failed', 'error': str(e)}), 500
@app.route('/images/search', methods=['GET'])
@token_required
def images_search():
    """Search images via Unsplash (if key provided) or return empty list."""
    q = request.args.get('q', '').strip()
    source = request.args.get('source', 'unsplash')
    limit = int(request.args.get('limit', 12))

    results = []
    try:
        if source == 'unsplash':
            key = os.environ.get('UNSPLASH_ACCESS_KEY')
            if not key:
                return jsonify({'images': []}), 200
            url = 'https://api.unsplash.com/search/photos'
            resp = requests.get(
                url,
                params={'query': q or 'photo', 'per_page': limit},
                headers={'Authorization': f'Client-ID {key}'},
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
            for hit in data.get('results', []):
                results.append({
                    'id': hit.get('id'),
                    'thumb': hit.get('urls', {}).get('thumb'),
                    'small': hit.get('urls', {}).get('small'),
                    'full': hit.get('urls', {}).get('full'),
                    'author': hit.get('user', {}).get('name')
                })
        else:
            # Google SERP requires API key or scraping; not implemented here.
            results = []
    except Exception as e:
        print('Image search error:', e)
        results = []

    return jsonify({'images': results}), 200


@app.route('/presentations/<pres_id>/slides/<slide_id>/ai-generate', methods=['POST'])
@token_required
def ai_generate_slide(pres_id, slide_id):
    """Generate slide content from a prompt with deep analysis using Anthropic/Gemini.
    Returns structured title + bullets and updates the slide."""
    payload = request.json or {}
    prompt = payload.get('prompt', '').strip()

    data = read_data()
    pres = data.get('presentations', {}).get(pres_id)
    if not pres:
        return jsonify({'message': 'Presentation not found'}), 404
    if pres['owner'] != request.user:
        return jsonify({'message': 'Forbidden'}), 403

    slide = next((s for s in pres['slides'] if s['id'] == slide_id), None)
    if not slide:
        return jsonify({'message': 'Slide not found'}), 404

    try:
        if prompt:
            llm_result = call_llm_for_structured_content(prompt)
            if llm_result:
                slide['title'] = llm_result.get('title', slide.get('title', 'Generated Slide'))
                bullets = llm_result.get('bullets', [])
                slide['content'] = '\n'.join([f"• {b}" for b in bullets]) if bullets else prompt
            else:
                slide['content'] = prompt

        pres['updated_at'] = datetime.datetime.utcnow().isoformat()
        write_data(data)
        return jsonify({'slide': slide}), 200
    except Exception as e:
        return jsonify({'message': 'AI generation failed', 'error': str(e)}), 500


# ============== CONTENT GENERATION ==============

@app.route('/generate', methods=['POST'])
@token_required
def generate():
    """
    Generate a presentation from text input with deep analysis using LLM.
    Modes:
    - 'ai': User provides details, LLM analyzes deeply and creates structured slides
    - 'user': User manually writes slide content (no LLM)
    - 'upload': User uploads document (uses document text)
    """
    try:
        payload = request.json or {}
        mode = payload.get('mode', 'ai')
        slide_count = max(1, min(int(payload.get('slide_count', 5)), 15))  # limit to 15 for perf
        title = payload.get('title', 'Generated Presentation').strip()
        details = payload.get('text', '').strip()
        
        if not title:
            return jsonify({'message': 'Title is required'}), 400
        
        slides = []
        
        if mode == 'ai' and details:
            # Use LLM to deeply analyze the input and create structured slides
            for i in range(slide_count):
                slide_prompt = f"""Presentation: "{title}"
                Input content: {details}

                Create slide {i+1} of {slide_count} that deeply analyzes and structures this information.
                Provide actionable insights and reasoning for this section."""
                                
                llm_result = call_llm_for_structured_content(slide_prompt)
                if llm_result:
                    slide_title = llm_result.get('title', f'{title} - Slide {i+1}')
                    bullets = llm_result.get('bullets', [])
                    slide_content = '\n'.join([f"• {b}" for b in bullets]) if bullets else 'Generated content'
                else:
                    slide_title = f'{title} - Slide {i+1}'
                    slide_content = 'Content to be filled'
                
                slides.append({
                    'id': str(uuid.uuid4()),
                    'title': slide_title,
                    'content': slide_content,
                    'image': None,
                    'style': {
                        'titleFontSize': 32,
                        'contentFontSize': 18,
                        'fontColor': '#000000',
                        'backgroundColor': '#ffffff',
                        'backgroundImage': None,
                        'backgroundOpacity': 100,
                        'backgroundBlur': 0
                    }
                })
        else:
            # Fallback: simple slide creation without LLM
            sentences = [s.strip() for s in details.split('.') if s.strip()]
            for i in range(slide_count):
                s_idx = min(i, len(sentences) - 1) if sentences else 0
                content = sentences[s_idx] if sentences else f'Slide {i+1} content'
                
                slides.append({
                    'id': str(uuid.uuid4()),
                    'title': f'{title} - Slide {i+1}',
                    'content': content,
                    'image': None,
                    'style': {
                        'titleFontSize': 32,
                        'contentFontSize': 18,
                        'fontColor': '#000000',
                        'backgroundColor': '#ffffff',
                        'backgroundImage': None,
                        'backgroundOpacity': 100,
                        'backgroundBlur': 0
                    }
                })
        
        # Save as presentation
        pres_id = str(uuid.uuid4())
        data = read_data()
        data['presentations'][pres_id] = {
            'id': pres_id,
            'owner': request.user,
            'title': title,
            'slides': slides,
            'created_at': datetime.datetime.utcnow().isoformat(),
            'updated_at': datetime.datetime.utcnow().isoformat()
        }
        write_data(data)
        
        return jsonify({'presentation': data['presentations'][pres_id]}), 201
    except Exception as e:
        return jsonify({'message': 'Generation failed', 'error': str(e)}), 500


@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'}), 200


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
